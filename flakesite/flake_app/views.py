from django.http import Http404, HttpRequest, HttpResponse, HttpResponseRedirect, FileResponse
from django.shortcuts import get_object_or_404, redirect, render
from django.urls import reverse
from django.views import generic

from django.contrib.auth.mixins import LoginRequiredMixin, PermissionRequiredMixin
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.contrib.messages.views import SuccessMessageMixin

from rest_framework import permissions
from rest_framework import generics

from rest_framework.renderers import TemplateHTMLRenderer
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework.parsers import MultiPartParser

from rules.contrib.views import AutoPermissionRequiredMixin

from .permissions import IsOwnerOrReadOnly
from .serializers import FlakePolymorphicSerializer
from .models import Flake, Graphene, hBN, Device, Comment
from .utils import get_contour, lookup_tables
from .forms import CommentForm

from pptx import Presentation
from pptx.shapes.freeform import FreeformBuilder
from pptx.util import Inches

from PIL import Image, ImageDraw, ImageEnhance
import cv2 as cv
import numpy as np

import io

class FlakeIndex(LoginRequiredMixin, generic.ListView):
    template_name = 'flake_app/flake-index.html'
    context_object_name = 'all_flakes'

    model = Flake
    def get_queryset(self):
        return Flake.objects.order_by('-box')

    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        all_graphenes = Graphene.objects.all()

        context['graphenes'] = all_graphenes
        # We pass a model instance to retreive the additional rows for the table programmatically.
        if len(all_graphenes):
            context['graphene_model'] = Graphene.objects.all()[0]

        context['hbns'] = hBN.objects.all()

        return context

def flake_image(request, pk):
    flake = get_object_or_404(Flake, pk = pk)
    if flake.flake_image.name.endswith('.dsx'):
        # Conversion of DSX files to JPEG.
        flake_image = Image.open(flake.flake_image)
        flake_image = flake_image.convert("RGB")
    
        if flake.contour:
            flake_arr = np.array(flake_image)
            contour = get_contour(flake.contour)
            cv.drawContours(flake_arr, [contour], 0, (255, 0, 0), 1)
            flake_image = Image.fromarray(flake_arr)

        response = HttpResponse(content_type = "image/jpg")
        flake_image.save(response, "JPEG")
        return response

    if flake.contour:
        # Convert to OpenCV formatted image
        flake_image = np.array(Image.open(flake.flake_image))
        
        contour = get_contour(flake.contour)

        cv.drawContours(flake_image, [contour], 0, (255, 0, 0), 1)

        flake_image = Image.fromarray(flake_image)
        response = HttpResponse(content_type = "image/jpg")
        flake_image.save(response, "JPEG")
        return response

    response = HttpResponse(flake.flake_image, content_type = "image/jpg")
    return response


def LUT_flake_image(request, pk):
    flake = get_object_or_404(Flake, pk = pk)
    flake_image = np.array(Image.open(flake.flake_image))

    R, G, B = cv.split(flake_image)

    adj_images = []

    for LUT in lookup_tables:
        R_LUT = cv.LUT(R, LUT[1])
        G_LUT = cv.LUT(G, LUT[2])
        B_LUT = cv.LUT(B, LUT[3])

        adj_images.append(cv.merge([R_LUT, G_LUT, B_LUT]))

    comb_image = np.hstack(adj_images)
    ret_image = Image.fromarray((255*comb_image).astype(np.uint8))
    response = HttpResponse(content_type = "image/jpg")

    ret_image.save(response, "JPEG")
    return response 

class FlakeDetail(LoginRequiredMixin, AutoPermissionRequiredMixin, generic.DetailView):
    template_name = 'flake_app/flake-detail.html'
    context_object_name = 'flake'
    queryset = Flake.objects.all()

    def post(self, request, pk, *args, **kwargs):
        # Adding the flake to the current device of the user.
        if 'add-to-device' in request.POST:
            if request.user.current_device:
                current_device = request.user.current_device
                flake_instance = Flake.objects.get(pk = pk)

                if not flake_instance.device and request.user.has_perm(Flake.get_perm("change"), flake_instance) and request.user.has_perm(Device.get_perm("change"), current_device):
                    flake_instance.device = current_device
                    flake_instance.save()
                return redirect(reverse('flake_app:flake-detail', args = (pk,)))

# TODO: Move this from a restframework view to an ordinary form. 
class FlakeEdit(LoginRequiredMixin, APIView):
    renderer_classes = [TemplateHTMLRenderer]
    parser_classes = [MultiPartParser]
    template_name = 'flake_app/flake-edit.html'
    permission_classes = [IsOwnerOrReadOnly]
    style = {'template_pack': 'rest_framework/vertical/'}

    model = Flake

    def get(self, request, pk):
        flake = get_object_or_404(Flake, pk = pk)
        poly_serializer = FlakePolymorphicSerializer(flake, context = {'request': request})
        return Response({'serializer': poly_serializer, 'flake' : flake, 'style' : self.style})
    
    def post(self, request, pk):
        flake = get_object_or_404(Flake, pk = pk)

        # Update the request to include the resourcetype for mapping to the actual serializer.
        updated_request = request.data.copy()
        updated_request.update({'resourcetype' : flake._meta.object_name})
        poly_serializer = FlakePolymorphicSerializer(flake, data = updated_request, context = {'request': request})
        if not poly_serializer.is_valid():
            return Response({'serializer': poly_serializer, 'flake': flake, 'style' : self.style})
        
        poly_serializer.save()
        # Redirect to the view page on completion.
        return redirect(reverse('flake_app:flake-detail', args = (pk,)))

@login_required()
def flake_map_image(request, pk):
    flake = get_object_or_404(Flake, pk = pk)
    map_image = Image.open(flake.map_image)
    map_draw = ImageDraw.Draw(map_image)
    map_draw.ellipse([flake.x_pos - 50, flake.y_pos - 50, flake.x_pos + 50, flake.y_pos + 50], fill = (255, 0, 0))

    response = HttpResponse(content_type = "image/jpg")
    map_image = map_image.resize((800, 800))
    map_image.save(response, "JPEG")
    return response

@login_required()
def flake_outlined_image(request, pk):
    flake = get_object_or_404(Flake, pk = pk)

    flake_image = np.array(Image.open(flake.flake_image))

    contour = np.frombuffer(flake.contour, dtype=np.int32).reshape(-1, 2)
    
    cv.drawContours(flake_image, [contour], 0, (0, 255, 0), 3)

    img_pil = Image.fromarray(flake_image)
    response = HttpResponse(content_type = 'image/jpg')
    img_pil.save(response, "JPEG")
    return response

class FlakeUpload(generics.CreateAPIView):
    queryset = Flake.objects.all()
    serializer_class = FlakePolymorphicSerializer
    permission_classes = [permissions.IsAuthenticated, IsOwnerOrReadOnly]

class FlakeUpdateRetrieve(generics.RetrieveUpdateAPIView):
    queryset = Flake.objects.all()
    serializer_class = FlakePolymorphicSerializer
    permission_classes = [permissions.IsAuthenticated, IsOwnerOrReadOnly]

class DeviceIndex(LoginRequiredMixin, generic.ListView):
    template_name = 'flake_app/device-index.html'
    context_object_name = 'device'

    model = Device

class DeviceDetail(LoginRequiredMixin, AutoPermissionRequiredMixin, generic.DetailView, generic.edit.FormMixin):
    template_name = 'flake_app/device-detail.html'
    context_object_name = 'device'
    queryset = Device.objects.all()

    model = Device

    form_class = CommentForm

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)

        device = self.get_object()

        parent_comments = Comment.objects.filter(device = device, parent_comment = None)

        def get_children(comment, curr_indent):
            all_comments = {}
            all_comments[comment] = curr_indent

            curr_indent += 40
            for reply in comment.replies.all():
                all_comments.update(get_children(reply, curr_indent))

            return all_comments
        
        comments_context = {}
        for parent in parent_comments:
            comments_context.update(get_children(parent, 0))

        context['sorted_comments'] = comments_context

        return context

    def post(self, request, pk, *args, **kwargs):
        device_instance = Device.objects.get(pk = pk)

        if 'remove-flake' in request.POST:
            if request.user.has_perm(Device.get_perm("change"), device_instance):
                removed = Flake.objects.get(pk = request.POST['remove-flake'])
                # TODO: Verify that the flake is actually part of the device.
                device_instance.flakes.remove(removed)
            else:
                messages.warning(request, "You must be the owner of the device to remove the flake.")
        elif 'set-current' in request.POST:
            request.user.current_device = device_instance
            request.user.save()
        elif 'post-comment' in request.POST:
            comment_form = CommentForm(data = request.POST)
            if comment_form.is_valid():
                body = comment_form.cleaned_data['body']

                if 'parent_comment' in comment_form.cleaned_data:
                    parent_comment = comment_form.cleaned_data['parent_comment']
                else:
                    parent_comment = None

                new_comment = Comment(body = body, parent_comment = parent_comment, device = self.get_object(), user = self.request.user)
                new_comment.save()

        return redirect(reverse('flake_app:device-detail', args = (pk,)))

class DeviceEdit(LoginRequiredMixin, AutoPermissionRequiredMixin, generic.UpdateView):
    model = Device
    queryset = Device.objects.all()
    fields = ['name', 'desc']
    template_name = 'flake_app/device-edit.html'

    def form_valid(self, form):
        device_form = form.save(commit = False)
        device_form.owner = self.request.user
        device_form.save()
        return redirect(reverse('flake_app:device-detail', args = (self.object.id,)))

class DeviceCreate(LoginRequiredMixin, generic.CreateView):
    model = Device
    fields = ['name', 'desc']
    template_name = 'flake_app/device-create.html'

    def form_valid(self, form):
        device_form = form.save(commit = False)
        device_form.owner = self.request.user
        device_form.save()
        return redirect(reverse('flake_app:device-index'))

@login_required
def device_powerpoint(request, pk):
    device_instance = Device.objects.get(pk = pk)
    if device_instance is None:
        return Http404

    # Create a new PowerPoint presentation with information regarding the device.

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = device_instance.name
    subtitle.text = "Device created by: {}".format(device_instance.owner)

    # Get each flake in the device, and add a slide with its relevant info.
    for flake_instance in device_instance.flakes.all():
        flake_slide_layout = prs.slide_layouts[6]
        
        flake_slide = prs.slides.add_slide(flake_slide_layout)

        slide_shapes = flake_slide.shapes

        PIL_image = Image.open(flake_instance.flake_image)
        flake_image = np.array(PIL_image)
    
        textBox = slide_shapes.add_textbox(0, 0, Inches(5), Inches(1))
        tf = textBox.text_frame
        tf.text = 'Name: {}\nBox: {}\nChip: {}'.format(flake_instance.name, flake_instance.box, flake_instance.chip)


        with io.BytesIO() as output:
            PIL_image.save(output, format = "png")
            image_width = Inches(flake_image.shape[1]/200)
            image_left = (prs.slide_width - image_width) / 2

            slide_shapes.add_picture(output, image_left, Inches(1), width = image_width)
        
        if flake_instance.contour:
            contour = np.frombuffer(flake_instance.contour, dtype=np.int32).reshape(-1, 2)
            approx = cv.approxPolyDP(contour, 0.001 * cv.arcLength(contour, True), True)

            # Reshape into array of tuples to be compatible with the pptx shapes
            reshaped_approx = approx.reshape(approx.shape[0], 2)

            form_builder = slide_shapes.build_freeform(reshaped_approx[0][0], reshaped_approx[0][1], Inches(1)/200)
            form_builder.add_line_segments(reshaped_approx[1 : ], close = True)
            free_shape = form_builder.convert_to_shape(image_left, Inches(1))
            free_shape.fill.background()
            free_shape.shadow.inherit = False

    # Create a buffer to hold the powerpoint object.
    buffer = io.BytesIO()
    
    prs.save(buffer)

    buffer.seek(0)
    return FileResponse(buffer, as_attachment = True, filename = device_instance.name + '.pptx')


